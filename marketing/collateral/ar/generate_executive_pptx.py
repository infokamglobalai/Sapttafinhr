#!/usr/bin/env python3
"""Generate Arabic Saptta executive PowerPoint (RTL-friendly Arabic text)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Saptta-Executive-Showcase-Jun2026-AR.pptx"

BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDES = [
    {
        "title": "منصة واحدة للأفراد والمال",
        "subtitle": "صبّتا HRMS + Accounts — workforce ومحاسبة امتثال بتسجيل دخول واحد.",
        "bullets": ["HR + Finance في منتج واحد", "30 موظف مشمول", "Complete: 95$/شهر", "14 يوم تجربة مجانية"],
        "title_slide": True,
    },
    {
        "title": "انفصال HR والمالية يضيع أسابيع",
        "bullets": [
            "رواتب في أداة وتصدير آخر للمحاسبة",
            "إجازات لا تطابق الحضور",
            "مديرون يراسلون HR بلا تدقيق",
            "صبّتا: سجل واحد وخدمة ذاتية",
        ],
    },
    {
        "title": "صبّتا في لمحة",
        "bullets": [
            "HRMS — 59$/شهر — موارد، وقت، رواتب، أداء",
            "Accounts — 59$/شهر — فواتير، GST، Tally",
            "Complete — 95$/شهر — توفير 20%",
        ],
    },
    {
        "title": "من نخدم",
        "bullets": [
            "الهند: PF/ESI/TDS/Form 16",
            "الخليج: HR كامل + رواتب PIFSS و indemnity",
            "تقنية، تصنيع، تجزئة، خدمات",
        ],
    },
    {
        "title": "من التوظيف حتى نهاية الخدمة",
        "bullets": [
            "توظيف · انضمام · تشغيل · رواتب · أداء · مغادرة",
            "ATS وصفحات careers و AI",
            "لوحة للقادة: headcount والحضور اليوم",
        ],
    },
    {
        "title": "صبّتا Accounts",
        "bullets": [
            "فواتير، مشتريات، مخزون، بنك",
            "GST، دفتر، Tally XML",
            "Complete: قيود رواتب → مالية",
        ],
    },
    {
        "title": "لماذا منصة موحدة",
        "bullets": [
            "headcount موثوق",
            "نهاية شهر أسرع",
            "مديرون self-service",
            "50 موظف: 125$/شهر Complete",
        ],
    },
    {
        "title": "امتثال الهند",
        "bullets": [
            "PF · ESI · PT · TDS · Form 16 · gratuity",
            "مراجعة → تشغيل → موافقة → قسائم",
            "TDS مبسّط — CA يتحقق قبل 24Q",
        ],
    },
    {
        "title": "الخليج والكويت",
        "bullets": [
            "KW / AE / SA من التسجيل",
            "PIFSS، GOSI، indemnity، IBAN",
            "تحقق مع مستشار محلي قبل التقديم",
        ],
    },
    {
        "title": "نتائج 90 يوماً",
        "bullets": [
            "60% وقت أقل للرواتب",
            "قسائم في نفس اليوم",
            "صفر إعادة إدخال",
            "سجل تدقيق كامل",
        ],
    },
    {
        "title": "تسعير بسيط",
        "bullets": [
            "HRMS 59 · Accounts 59 · Complete 95",
            "الهند: ₹7,999 + GST",
            "أقل من ZenHR ~$150 HR-only",
        ],
    },
    {
        "title": "الخطوة التالية: عرض 45 دقيقة",
        "bullets": [
            "1: اكتشاف — headcount وأدوات",
            "2: عرض حي HR + Finance",
            "3: تجربة خلال 24 ساعة",
            "hello@saptta.com · saptta.com",
        ],
    },
]


def _set_para(para, *, size=16, color=MUTED, bold=False, rtl=True):
    para.font.size = Pt(size)
    para.font.color.rgb = color
    para.font.bold = bold
    para.font.name = "Arial"
    if rtl:
        try:
            para._element.get_or_add_pPr().attrib[
                "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rtl"
            ] = "1"
        except Exception:
            pass


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        is_title = spec.get("title_slide")

        if is_title:
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = DARK
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.2))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.text = spec["title"]
            _set_para(tf.paragraphs[0], size=36, color=WHITE, bold=True)
            sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(10), Inches(1.4))
            sub.text_frame.text = spec.get("subtitle", "")
            _set_para(sub.text_frame.paragraphs[0], size=18, color=RGBColor(0x93, 0xC5, 0xFD))
        else:
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
            bar.fill.solid()
            bar.fill.fore_color.rgb = BLUE
            bar.line.fill.background()
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11), Inches(1.2))
            title_box.text_frame.text = spec["title"]
            _set_para(title_box.text_frame.paragraphs[0], size=26, color=DARK, bold=True)
            body = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.5))
            tf = body.text_frame
            tf.clear()
            for i, item in enumerate(spec["bullets"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item
                _set_para(p, size=17, color=MUTED)

        foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(12), Inches(0.4))
        foot.text_frame.text = f"صبّتا · عرض تنفيذي · {idx}/12 · يونيو 2026 · EN: ../en/"
        _set_para(foot.text_frame.paragraphs[0], size=10, color=MUTED, rtl=False)

    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    build()
