#!/usr/bin/env python3
"""
Build executive CEO/manager platform deck with product screenshots,
one slide per AI feature, and speaker notes for live presentations.

Usage:
    python marketing/collateral/build_platform_ceo_deck_pptx.py

Output:
    marketing/collateral/Saptta-Platform-CEO-Deck.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
COLLATERAL = Path(__file__).resolve().parent
OUT = COLLATERAL / "Saptta-Platform-CEO-Deck.pptx"
GENERATED = COLLATERAL / "generated"
SHOTS = [
    ROOT / "shots" / "02-switcher.png",
    ROOT / "shots" / "04-hr.png",
    ROOT / "shots" / "03-finance.png",
]

# Brand palette
BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_DARK = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_SOFT = RGBColor(0xEF, 0xF6, 0xFF)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
GREEN_SOFT = RGBColor(0xEC, 0xFD, 0xF5)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Segoe UI"


def rgb(r, g, b) -> RGBColor:
    return RGBColor(r, g, b)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(shape, text: str, *, size=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def _notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def _footer(slide, idx: int, total: int):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), H - Inches(0.38), W, Inches(0.38))
    _fill(bar, LIGHT_BG)
    box = slide.shapes.add_textbox(Inches(0.6), H - Inches(0.32), Inches(12), Inches(0.28))
    _text(
        box,
        f"Saptta  ·  Executive Platform Brief  ·  {idx}/{total}  ·  saptta.com  ·  hello@saptta.com",
        size=9,
        color=SLATE,
    )


def _header(slide, title: str, tag: str = ""):
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.35), Inches(0.45), Inches(0.45))
    _fill(logo, BLUE)
    logo.adjustments[0] = 0.15
    lt = slide.shapes.add_textbox(Inches(0.62), Inches(0.42), Inches(0.35), Inches(0.35))
    _text(lt, "S", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    brand = slide.shapes.add_textbox(Inches(1.05), Inches(0.38), Inches(2), Inches(0.4))
    _text(brand, "Saptta", size=14, color=NAVY, bold=True)

    if tag:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(0.78), Inches(2.2), Inches(0.28))
        _fill(pill, BLUE_SOFT)
        pill.line.color.rgb = rgb(0xBF, 0xDB, 0xFE)
        pt = slide.shapes.add_textbox(Inches(1.12), Inches(0.8), Inches(2.1), Inches(0.24))
        _text(pt, tag.upper(), size=8, color=BLUE, bold=True)

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12), Inches(0.75))
    _text(title_box, title, size=26, color=NAVY, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.95), Inches(12.2), Inches(0.03))
    _fill(line, BORDER)


def _card(slide, left, top, width, height, title, lines, *, accent=BLUE, featured=False):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.08
    if featured:
        _fill(rect, BLUE_SOFT)
        rect.line.color.rgb = BLUE
        rect.line.width = Pt(2)
    else:
        _fill(rect, WHITE)
        rect.line.color.rgb = BORDER
        rect.line.width = Pt(1)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    _fill(stripe, accent)

    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.35), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = FONT
    p0.font.size = Pt(12 if not featured else 13)
    p0.font.bold = True
    p0.font.color.rgb = NAVY if not featured else BLUE_DARK
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(10)
        p.font.color.rgb = SLATE
        p.space_before = Pt(4)


def _stat(slide, left, top, value, label, *, primary=False):
    w, h = Inches(2.85), Inches(1.35)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.adjustments[0] = 0.12
    _fill(box, BLUE_SOFT if primary else LIGHT_BG)
    if primary:
        box.line.color.rgb = rgb(0xBF, 0xDB, 0xFE)
    else:
        box.line.color.rgb = BORDER

    v = slide.shapes.add_textbox(left, top + Inches(0.22), w, Inches(0.55))
    _text(v, value, size=26 if primary else 22, color=BLUE if primary else NAVY, bold=True, align=PP_ALIGN.CENTER)
    l = slide.shapes.add_textbox(left, top + Inches(0.78), w, Inches(0.4))
    _text(l, label, size=9, color=SLATE, align=PP_ALIGN.CENTER)


def _image(slide, path: Path, left, top, width, *, caption: str = ""):
    if path.is_file():
        pic = slide.shapes.add_picture(str(path), left, top, width=width)
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left - Inches(0.06),
            top - Inches(0.06),
            width + Inches(0.12),
            pic.height + Inches(0.12),
        )
        frame.adjustments[0] = 0.04
        _fill(frame, WHITE)
        frame.line.color.rgb = BORDER
        spTree = slide.shapes._spTree
        sp = pic._element
        spTree.remove(sp)
        spTree.insert(2, sp)
        img_h = pic.height
    else:
        placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(3.5))
        placeholder.adjustments[0] = 0.04
        _fill(placeholder, LIGHT_BG)
        placeholder.line.color.rgb = BORDER
        _text(
            slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(1)),
            f"[Image: {path.name}]",
            size=11,
            color=SLATE,
            align=PP_ALIGN.CENTER,
        )
        img_h = Inches(3.5)
    if caption:
        cap = slide.shapes.add_textbox(left, top + img_h + Inches(0.08), width, Inches(0.35))
        _text(cap, caption, size=10, color=SLATE, bold=True)


def _bullets(slide, left, top, width, height, items, *, title: str = ""):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            label, body = item
            p.text = label
            p.font.name = FONT
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = NAVY
            p2 = tf.add_paragraph()
            p2.text = body
            p2.font.name = FONT
            p2.font.size = Pt(11)
            p2.font.color.rgb = SLATE
            p2.space_after = Pt(8)
        else:
            p.text = f"•  {item}"
            p.font.name = FONT
            p.font.size = Pt(12)
            p.font.color.rgb = SLATE
            p.space_after = Pt(8)
    if title:
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = NAVY


def slide_title(prs, total):
    s = _blank(prs)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    orb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.2), Inches(5), Inches(5))
    _fill(orb, BLUE_DARK)
    orb.line.fill.background()

    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.95), Inches(3.2), Inches(0.35))
    _fill(badge, rgb(0x1E, 0x3A, 0x8A))
    b = s.shapes.add_textbox(Inches(1.0), Inches(0.98), Inches(3.0), Inches(0.3))
    _text(b, "EXECUTIVE PLATFORM BRIEF · JUN 2026", size=9, color=rgb(0x93, 0xC5, 0xFD), bold=True)

    t = s.shapes.add_textbox(Inches(0.85), Inches(1.5), Inches(10.5), Inches(1.8))
    _text(t, "Saptta Complete:\nPeople, books & AI", size=42, color=WHITE, bold=True)

    sub = s.shapes.add_textbox(Inches(0.88), Inches(3.35), Inches(9.5), Inches(1.2))
    _text(
        sub,
        "One Indian-ready platform for HRMS, finance, compliance, and\nscoped AI assistants — built for CEOs, CHROs, and CFOs.",
        size=17,
        color=rgb(0xCB, 0xD5, 0xE1),
    )

    stats_y = Inches(4.65)
    for i, (val, lbl) in enumerate([
        ("2-in-1", "HR + Finance"),
        ("10", "AI features"),
        ("₹7,999", "Complete / mo"),
        ("30", "Employees incl."),
    ]):
        _stat(s, Inches(0.85) + i * Inches(3.05), stats_y, val, lbl, primary=(i == 2))

    _footer(s, 1, total)
    _notes(
        s,
        "Open with the CEO headline: you are not buying another HR tool or another accounting tool — "
        "you are buying one operating system for people and money. Mention this deck covers the full platform, "
        "each AI capability, and what your leadership team gets in the first 90 days.",
    )
    return s


def slide_executive_summary(prs, total, idx):
    s = _blank(prs)
    _header(s, "Why leaders choose Saptta", "Executive summary")
    _card(
        s, Inches(0.55), Inches(2.15), Inches(3.9), Inches(4.5),
        "For the CEO",
        [
            "One login, one employee truth",
            "Payroll posts to ledger automatically",
            "Role-based access — no data leaks",
            "AI answers from your data, not the internet",
        ],
        accent=NAVY,
        featured=True,
    )
    _card(
        s, Inches(4.65), Inches(2.15), Inches(3.9), Inches(4.5),
        "For the CHRO",
        [
            "Hire → onboard → pay → exit lifecycle",
            "Manager self-service for leave & attendance",
            "Recruitment AI: JD, rank, offer letter",
            "Policy Q&A from your uploaded docs",
        ],
        accent=BLUE,
    )
    _card(
        s, Inches(8.75), Inches(2.15), Inches(3.9), Inches(4.5),
        "For the CFO",
        [
            "GST invoices, ledger, bank reconciliation",
            "Finance AI for P&L, receivables, GST",
            "AI bank reconciliation suggestions",
            "Tally export + SaaS billing invoice",
        ],
        accent=GREEN,
    )
    _footer(s, idx, total)
    _notes(
        s,
        "Frame the three personas in the room. CEO cares about unified truth and risk. CHRO cares about "
        "manager time saved and compliance. CFO cares about books matching payroll without re-keying. "
        "Saptta Complete is the bundle that connects all three.",
    )


def slide_problem(prs, total, idx):
    s = _blank(prs)
    _header(s, "The hidden cost of disconnected systems", "Challenge")
    _card(
        s, Inches(0.55), Inches(2.15), Inches(5.9), Inches(3.8),
        "Without a unified platform",
        [
            "✗  HR and Finance run on different tools",
            "✗  Payroll re-keyed into accounting",
            "✗  Managers email HR for every approval",
            "✗  No audit trail on sensitive changes",
            "✗  Generic ChatGPT cannot read your books",
        ],
        accent=ORANGE,
    )
    _card(
        s, Inches(6.85), Inches(2.15), Inches(5.9), Inches(3.8),
        "With Saptta Complete",
        [
            "✓  One employee record — one truth",
            "✓  Attendance → payroll → payslip → ledger",
            "✓  Manager & employee self-service",
            "✓  PF, GST, Tally exports built-in",
            "✓  Scoped AI on tenant data only",
        ],
        accent=GREEN,
        featured=True,
    )
    callout = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.55))
    callout.adjustments[0] = 0.2
    _fill(callout, BLUE_SOFT)
    callout.line.color.rgb = rgb(0xBF, 0xDB, 0xFE)
    c = s.shapes.add_textbox(Inches(0.75), Inches(6.12), Inches(11.8), Inches(0.45))
    _text(c, "CEO insight: leadership time spent reconciling two versions of your company is the real cost.", size=11, color=BLUE_DARK, bold=True)
    _footer(s, idx, total)
    _notes(
        s,
        "Use a concrete example: finance asks HR for headcount cost, HR sends a spreadsheet, finance re-enters "
        "into Tally. That loop disappears when payroll journals post automatically.",
    )


def slide_platform(prs, total, idx):
    s = _blank(prs)
    _header(s, "Three products. One platform.", "Platform")
    plans = [
        ("Saptta HRMS", "₹4,999/mo", ["Payroll & statutory", "Leave & attendance", "Recruitment & performance", "HR AI assistants"], BLUE, False),
        ("Saptta Finance", "₹4,999/mo", ["Invoicing & GST", "Ledger & bank rec", "Inventory & purchases", "Finance AI assistant"], PURPLE, False),
        ("Saptta Complete ★", "₹7,999/mo", ["Everything in HRMS + Finance", "Save ₹1,999 vs separate", "Payroll → ledger sync", "Recommended for owners"], BLUE, True),
    ]
    x0 = Inches(0.55)
    for i, (name, price, feats, accent, feat) in enumerate(plans):
        left = x0 + i * Inches(4.15)
        _card(s, left, Inches(2.2), Inches(3.95), Inches(3.5), name, [price, ""] + feats, accent=accent, featured=feat)
    _image(s, GENERATED / "saptta-unified-platform-hrms-accounts.png", Inches(0.55), Inches(5.85), Inches(12.2))
    _footer(s, idx, total)
    _notes(
        s,
        "Explain pricing is ex-GST. HRMS and Complete include 30 employees (+₹111/extra). Finance has unlimited "
        "users. Complete is the strategic choice because payroll-to-ledger sync only exists in the bundle.",
    )


def slide_screenshot(prs, total, idx, title, tag, img_path, caption, bullets, notes):
    s = _blank(prs)
    _header(s, title, tag)
    _image(s, img_path, Inches(6.6), Inches(2.15), Inches(6.2), caption=caption)
    _bullets(s, Inches(0.55), Inches(2.2), Inches(5.8), Inches(4.2), bullets)
    _footer(s, idx, total)
    _notes(s, notes)


def slide_payroll_sync(prs, total, idx):
    s = _blank(prs)
    _header(s, "When HR runs payroll, finance updates itself", "Complete bundle")
    _image(
        s,
        GENERATED / "saptta-payroll-ledger-sync.png",
        Inches(6.4),
        Inches(2.15),
        Inches(6.4),
        caption="Payroll → journal → ledger → GSTR-ready",
    )
    _bullets(
        s,
        Inches(0.55),
        Inches(2.2),
        Inches(5.7),
        Inches(4.5),
        [
            ("Automatic journal entries", "Salary, PF, ESI, and reimbursements post with correct debit/credit — no CSV bridge."),
            ("Single employee master", "Headcount cost in finance matches HR payroll runs."),
            ("Audit trail", "Every payroll run and journal is traceable."),
            ("CFO benefit", "Month-end close starts with books already updated."),
        ],
    )
    _footer(s, idx, total)
    _notes(
        s,
        "This is the #1 differentiator vs Keka-only or Tally-only stacks. Walk through: HR approves payroll → "
        "system generates payslips AND posts journals → finance sees updated ledger → GST reports stay aligned.",
    )


def slide_rbac(prs, total, idx):
    s = _blank(prs)
    _header(s, "Who logs in — roles your team understands", "Security & RBAC")
    _card(
        s, Inches(0.55), Inches(2.15), Inches(5.9), Inches(2.0),
        "After purchase",
        ["Signup email = workspace owner (full admin)", "Invite HR admin, managers, employees", "Finance: Owner → Admin → Accountant → Viewer"],
        accent=BLUE,
        featured=True,
    )
    _card(
        s, Inches(0.55), Inches(4.35), Inches(5.9), Inches(2.0),
        "Not sold to customers",
        ["Saptta platform superadmin = Saptta staff only", "Manages billing, support, all tenants", "Buyers never see /superadmin"],
        accent=ORANGE,
    )
    _card(
        s, Inches(6.85), Inches(2.15), Inches(5.9), Inches(4.2),
        "Granular permissions (HR)",
        [
            "HR admin: full HR operations",
            "Manager: team leave, attendance, reviews",
            "Employee: self-service ESS only",
            "Finance users: separate role matrix",
            "Every route checks permission — not just menu hiding",
        ],
        accent=GREEN,
        featured=True,
    )
    _footer(s, idx, total)
    _notes(
        s,
        "CEOs worry about data leakage. Explain workspace owner is the purchaser — not Saptta staff. "
        "Managers see only their team. Employees see only themselves. Finance roles are independent.",
    )


def slide_ai_divider(prs, total, idx):
    s = _blank(prs)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BLUE_DARK

    t = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11), Inches(1.5))
    _text(t, "AI built for business —\nnot generic chat", size=40, color=WHITE, bold=True)

    sub = s.shapes.add_textbox(Inches(0.92), Inches(3.6), Inches(10), Inches(1.2))
    _text(
        sub,
        "10 scoped AI capabilities · Tenant data only · Each assistant stays in its lane\n"
        "Refuses legal advice, cross-module queries, and data outside your company",
        size=16,
        color=rgb(0xBF, 0xDB, 0xFE),
    )
    _footer(s, idx, total)
    _notes(
        s,
        "Transition slide. Key message: this is NOT ChatGPT pasted into the product. Each assistant uses "
        "Claude with tool-calling against live database queries scoped to the tenant. No hallucinated headcounts.",
    )


def slide_ai_architecture(prs, total, idx):
    s = _blank(prs)
    _header(s, "Three assistants — each knows its lane", "AI architecture")
    _image(s, GENERATED / "saptta-ai-three-assistants.png", Inches(0.55), Inches(2.15), Inches(6.5), caption="HR · Finance · Saptta Guide")
    _card(
        s, Inches(7.3), Inches(2.15), Inches(5.45), Inches(1.35),
        "HR Assistant",
        ["Leave, attendance, payroll, team summaries", "Refuses finance questions → handoff"],
        accent=BLUE,
    )
    _card(
        s, Inches(7.3), Inches(3.65), Inches(5.45), Inches(1.35),
        "Finance Assistant",
        ["Invoices, P&L, GST, receivables", "Refuses HR questions → handoff"],
        accent=GREEN,
    )
    _card(
        s, Inches(7.3), Inches(5.15), Inches(5.45), Inches(1.35),
        "Saptta Guide",
        ["Product & onboarding Q&A", "No live tenant data — routes to HR/Finance AI"],
        accent=PURPLE,
    )
    _footer(s, idx, total)
    _notes(
        s,
        "Technical how-it-works for the CTO in the room: Claude API + tool use. The model decides which "
        "database query to run (get_leave_balance, get_gst_summary, etc.) and only speaks from tool results.",
    )


def slide_ai_feature(
    prs,
    total,
    idx,
    *,
    title: str,
    tag: str,
    who: str,
    bullets: list,
    how_it_works: list,
    example_prompt: str,
    image: Path | None,
    notes: str,
    accent=BLUE,
):
    s = _blank(prs)
    _header(s, title, tag)

    who_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(2.1), Inches(12.2), Inches(0.42))
    who_box.adjustments[0] = 0.3
    _fill(who_box, BLUE_SOFT)
    who_box.line.color.rgb = rgb(0xBF, 0xDB, 0xFE)
    wt = s.shapes.add_textbox(Inches(0.75), Inches(2.16), Inches(11.8), Inches(0.35))
    _text(wt, f"Who uses it: {who}", size=11, color=BLUE_DARK, bold=True)

    if image and image.is_file():
        _image(s, image, Inches(7.0), Inches(2.65), Inches(5.8))
        content_left = Inches(0.55)
        content_width = Inches(6.2)
    else:
        content_left = Inches(0.55)
        content_width = Inches(12.2)

    _card(
        s, content_left, Inches(2.65), content_width, Inches(2.1),
        "What it does",
        bullets,
        accent=accent,
    )
    _card(
        s, content_left, Inches(4.9), content_width, Inches(1.55),
        "How it works (under the hood)",
        how_it_works,
        accent=NAVY,
    )

    prompt_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.48))
    prompt_box.adjustments[0] = 0.2
    _fill(prompt_box, GREEN_SOFT)
    prompt_box.line.color.rgb = rgb(0xA7, 0xF3, 0xD0)
    pt = s.shapes.add_textbox(Inches(0.75), Inches(6.62), Inches(11.8), Inches(0.4))
    _text(pt, f'Example: "{example_prompt}"', size=11, color=GREEN, bold=True)

    _footer(s, idx, total)
    _notes(s, notes)


def slide_recruitment_overview(prs, total, idx):
    s = _blank(prs)
    _header(s, "Recruitment AI — hire faster with less bias", "Recruitment AI")
    _image(s, GENERATED / "saptta-recruitment-ai-flow.png", Inches(0.55), Inches(2.15), Inches(6.3), caption="JD → Parse → Rank → Offer")
    _card(
        s, Inches(7.1), Inches(2.15), Inches(5.65), Inches(1.5),
        "1. JD Generator",
        ["Role title + skills → professional JD in seconds", "Structured JSON for careers page"],
        accent=BLUE,
    )
    _card(
        s, Inches(7.1), Inches(3.8), Inches(5.65), Inches(1.5),
        "2. Resume Parse & Rank",
        ["Extract skills, experience, education from PDF", "Score candidates against job requirements"],
        accent=PURPLE,
    )
    _card(
        s, Inches(7.1), Inches(5.45), Inches(5.65), Inches(1.5),
        "3. Offer Letter",
        ["Generate offer from approved template", "HR reviews before sending"],
        accent=GREEN,
    )
    _footer(s, idx, total)
    _notes(
        s,
        "Overview before detail slides. HR team saves days per hire. AI ranks consistently — human always "
        "makes final decision. Next three slides go deeper on each step.",
    )


def slide_modules(prs, total, idx):
    s = _blank(prs)
    _header(s, "Hire to retire — full HR lifecycle", "HRMS")
    modules = [
        ("Recruit", "ATS, careers, AI rank"),
        ("Onboard", "Checklists, documents"),
        ("Operate", "Attendance, leave"),
        ("Pay", "Run, payslip, bank"),
        ("Grow", "Performance reviews"),
        ("Comply", "Letters, policies"),
        ("Engage", "ESS portal"),
        ("Exit", "Settlement, revoke"),
    ]
    for i, (name, desc) in enumerate(modules):
        col, row = i % 4, i // 4
        left = Inches(0.55) + col * Inches(3.1)
        top = Inches(2.15) + row * Inches(1.85)
        _card(s, left, top, Inches(2.85), Inches(1.55), name, [desc], accent=BLUE)
    _footer(s, idx, total)
    _notes(s, "Quick module map for CHRO. Emphasize everything connects — attendance feeds payroll feeds compliance.")


def slide_compliance(prs, total, idx):
    s = _blank(prs)
    _header(s, "Compliance built for real businesses", "Trust")
    _image(s, GENERATED / "saptta-compliance-ai.png", Inches(8.0), Inches(2.15), Inches(4.75))
    _card(
        s, Inches(0.55), Inches(2.15), Inches(7.2), Inches(2.1),
        "India (production-ready)",
        ["PF ECR · ESI · PT · TDS · Form 16", "Gratuity accrual + exit estimate", "Salary register · Bank advice · Tally XML"],
        accent=GREEN,
        featured=True,
    )
    _card(
        s, Inches(0.55), Inches(4.45), Inches(7.2), Inches(2.1),
        "GCC & Kuwait (demo-ready MVP)",
        ["KW · AE · SA jurisdiction signup", "PIFSS · GOSI · indemnity accrual", "ZATCA · Peppol integration path"],
        accent=ORANGE,
    )
    _footer(s, idx, total)
    _notes(s, "Be honest on GCC: demo-ready MVP, validate rates with local consultant. India is production-ready core.")


def slide_pricing(prs, total, idx):
    s = _blank(prs)
    _header(s, "Simple pricing — one decision for leadership", "Pricing")
    rows = [
        ["Plan", "INR / mo (ex-GST)", "Includes", "Best for"],
        ["HRMS", "₹4,999", "30 employees, HR AI", "HR-led companies"],
        ["Finance", "₹4,999", "Unlimited finance users", "CA / accounts-led"],
        ["Complete ★", "₹7,999", "HR + Finance + sync + all AI", "Owners wanting one platform"],
    ]
    tbl = s.shapes.add_table(4, 4, Inches(0.55), Inches(2.15), Inches(12.2), Inches(2.2)).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.name = FONT
                para.font.size = Pt(10 if r else 9)
                para.font.bold = r == 0 or (r == 3 and c == 0)
                para.font.color.rgb = NAVY if r == 3 else (SLATE if r else WHITE)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif r == 3:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_SOFT

    for i, (val, lbl) in enumerate([("50 emp", "→ ₹8,219"), ("100 emp", "→ ₹13,769"), ("+30", "+₹111/emp")]):
        _stat(s, Inches(0.55) + i * Inches(4.05), Inches(4.65), val, lbl, primary=(i == 0))

    callout = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.45))
    callout.adjustments[0] = 0.2
    _fill(callout, GREEN_SOFT)
    c = s.shapes.add_textbox(Inches(0.75), Inches(6.22), Inches(11.8), Inches(0.35))
    _text(c, "Complete saves ₹1,999/mo vs buying HRMS + Finance separately — plus payroll-to-ledger sync.", size=11, color=GREEN, bold=True)
    _footer(s, idx, total)
    _notes(s, "USD equivalent ~$59/$59/$95 for international buyers. GST invoice emailed on activation.")


def slide_outcomes(prs, total, idx):
    s = _blank(prs)
    _header(s, "Results leadership sees in 90 days", "Outcomes")
    for i, (val, lbl) in enumerate([
        ("60%", "Less payroll prep"),
        ("Same day", "Payslip publish"),
        ("Zero", "Spreadsheet re-key"),
        ("100%", "Audit on changes"),
    ]):
        _stat(s, Inches(0.55) + i * Inches(3.15), Inches(2.3), val, lbl, primary=(i == 0))
    _card(s, Inches(0.55), Inches(4.0), Inches(5.9), Inches(2.3), "CHRO wins", ["Manager leave approvals in-app", "Recruitment pipeline + AI ranking", "Policy Q&A reduces HR inbox"], accent=BLUE)
    _card(s, Inches(6.85), Inches(4.0), Inches(5.9), Inches(2.3), "CFO wins", ["Payroll cost visible in finance", "AI bank reconcile suggestions", "GST + ledger same platform"], accent=GREEN)
    _footer(s, idx, total)
    _notes(s, "Use customer language: less time, fewer errors, faster close. Offer a 45-min demo on their scenarios.")


def slide_cta(prs, total, idx):
    s = _blank(prs)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    t = s.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(11), Inches(1.2))
    _text(t, "See your company on Saptta", size=36, color=WHITE, bold=True)

    sub = s.shapes.add_textbox(Inches(0.92), Inches(2.35), Inches(10), Inches(0.6))
    _text(sub, "45-minute live demo on your scenarios · Workspace live after payment", size=16, color=rgb(0x93, 0xC5, 0xFD))

    steps = [
        ("1", "Discovery", "Headcount, tools, pain points"),
        ("2", "Live demo", "HR + Finance + AI, one login"),
        ("3", "Go live", "Invite team, configure roles"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        left = Inches(0.9) + i * Inches(4.05)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.2), Inches(3.7), Inches(1.65))
        box.adjustments[0] = 0.1
        _fill(box, rgb(0x1E, 0x29, 0x3B))
        box.line.color.rgb = rgb(0x33, 0x41, 0x55)
        n = s.shapes.add_textbox(left + Inches(0.2), Inches(3.35), Inches(0.5), Inches(0.4))
        _text(n, num, size=20, color=BLUE, bold=True)
        tt = s.shapes.add_textbox(left + Inches(0.55), Inches(3.32), Inches(2.9), Inches(0.4))
        _text(tt, title, size=14, color=WHITE, bold=True)
        dd = s.shapes.add_textbox(left + Inches(0.2), Inches(3.85), Inches(3.3), Inches(0.7))
        _text(dd, desc, size=11, color=rgb(0x94, 0xA3, 0xB8))

    contact = s.shapes.add_textbox(Inches(0.9), Inches(5.35), Inches(11), Inches(0.8))
    _text(contact, "hello@saptta.com  ·  saptta.com  ·  GST invoice on activation  ·  Free Complete setup (launch offer)", size=14, color=WHITE, align=PP_ALIGN.CENTER)

    _footer(s, idx, total)
    _notes(s, "Close with next step: book demo or start signup. Offer to walk through RBAC setup for their org chart.")


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_spec = []  # track for total count
    # We build incrementally; count at end
    idx = 0

    def next_idx():
        nonlocal idx
        idx += 1
        return idx

    # Pre-count: title + exec + problem + platform + 3 screenshots + payroll + rbac + ai divider +
    # ai arch + 8 ai features + recruitment overview + modules + compliance + pricing + outcomes + cta
    total = 27

    slide_title(prs, total)
    slide_executive_summary(prs, total, next_idx())
    slide_problem(prs, total, next_idx())
    slide_platform(prs, total, next_idx())

    slide_screenshot(
        prs, total, next_idx(),
        "One login. Two products.",
        "Live product",
        SHOTS[0],
        "Product switcher — HR & Finance",
        [
            "Unified platform login at your workspace URL",
            "Switch HR ↔ Finance instantly — no re-login",
            "Same users, same employee master",
            "Complete plan unlocks payroll-to-ledger sync",
        ],
        "Show the switcher live. CEO should see this is one vendor, one contract, one support channel.",
    )
    slide_screenshot(
        prs, total, next_idx(),
        "HR dashboard built for leaders",
        "Saptta HRMS",
        SHOTS[1],
        "Headcount, attendance, priorities",
        [
            "Today's present / absent at a glance",
            "Pending leave & payroll actions surfaced",
            "Manager and employee self-service (ESS)",
            "HR AI chat in every screen",
        ],
        "CHRO slide. Point out pending actions — this replaces morning email chains.",
    )
    slide_screenshot(
        prs, total, next_idx(),
        "Finance your CA will use",
        "Saptta Accounts",
        SHOTS[2],
        "Invoicing, ledger, GST",
        [
            "Full double-entry accounting workspace",
            "GST-ready: GSTR-1/3B, e-invoice path",
            "Tally XML export for your CA",
            "Finance AI for live P&L and receivables",
        ],
        "CFO slide. Emphasize unlimited finance users on flat fee — no per-seat tax.",
    )

    slide_payroll_sync(prs, total, next_idx())
    slide_rbac(prs, total, next_idx())
    slide_ai_divider(prs, total, next_idx())
    slide_ai_architecture(prs, total, next_idx())

    # --- Individual AI feature slides ---
    slide_ai_feature(
        prs, total, next_idx(),
        title="HR Employee Assistant — self-service in plain English",
        tag="AI · Employee ESS",
        who="Every employee with ESS access",
        bullets=[
            "Ask leave balance, upcoming holidays, attendance summary",
            "Apply leave or submit expense — AI confirms before saving",
            "View payslip-related answers scoped to the logged-in employee",
        ],
        how_it_works=[
            "Claude API with tool-calling (get_leave_balance, apply_leave, etc.)",
            "Tools query only the current tenant's HR database",
            "Confirms destructive actions (leave apply) before executing",
        ],
        example_prompt="What is my casual leave balance? Apply casual leave next Friday.",
        image=GENERATED / "saptta-employee-ess-mobile.png",
        notes=(
            "Demo prompt live. Show that employee never needs to navigate 5 menus. "
            "AI refuses to show another employee's salary. Manager gets team summary variant."
        ),
        accent=BLUE,
    )

    slide_ai_feature(
        prs, total, next_idx(),
        title="HR Management Assistant — answers for HR admins",
        tag="AI · HR Admin",
        who="HR admins and workspace owners",
        bullets=[
            "Headcount, attendance rates, pending leave approvals",
            "Payroll run status and summary for the organisation",
            "Team-on-leave queries for workforce planning",
        ],
        how_it_works=[
            "Separate admin-scoped tools (org-wide read access)",
            "Same tenant isolation — never crosses companies",
            "Refuses finance questions → directs to Finance AI",
        ],
        example_prompt="How many employees are present today? Any pending leave approvals?",
        image=GENERATED / "saptta-hr-hero-dashboard.png",
        notes="For CHRO/HR head. Contrast with employee assistant — broader read scope, still HR-only.",
        accent=BLUE,
    )

    slide_ai_feature(
        prs, total, next_idx(),
        title="Policy Q&A — answers from your uploaded policies",
        tag="AI · Policy Hub",
        who="Employees (via HR Assistant) and HR admins",
        bullets=[
            "Answers ONLY from tenant-uploaded policy documents",
            "Cites policy title — does not invent rules",
            "Reduces repetitive HR inbox questions",
        ],
        how_it_works=[
            "Loads active PolicyDocument records (up to 12k chars context)",
            "Claude reads policy text + employee question",
            "If answer not in docs → says so and suggests contacting HR",
        ],
        example_prompt="What is our work-from-home policy for Mumbai office?",
        image=GENERATED / "saptta-compliance-ai.png",
        notes="Upload PDFs/text under HR → Policies. Great for onboarding 30+ employee companies.",
        accent=PURPLE,
    )

    slide_ai_feature(
        prs, total, next_idx(),
        title="Finance AI Assistant — your books in conversation",
        tag="AI · Finance",
        who="Finance Owner, Admin, Accountant roles",
        bullets=[
            "Outstanding invoices, cash position, P&L, GST liability",
            "Draft payment reminder emails for overdue customers",
            "Trial balance and receivables without running reports manually",
        ],
        how_it_works=[
            "Tool-calling: get_overdue_invoices, get_gst_summary, get_pl, etc.",
            "Strict refusal of HR/payroll topics → HR Assistant handoff",
            "Amounts quoted in ₹ from live ledger — never guessed",
        ],
        example_prompt="Show overdue invoices above ₹50,000. What is our GST liability this month?",
        image=GENERATED / "saptta-accounts-finance-dashboard.png",
        notes="CFO favorite. Demo with real overdue list. Stress it won't answer 'who is on leave' — by design.",
        accent=GREEN,
    )

    slide_ai_feature(
        prs, total, next_idx(),
        title="Bank AI Reconciliation — match statements faster",
        tag="AI · Finance",
        who="Accountants and finance admins",
        bullets=[
            "Suggests matches between bank statement lines and ledger entries",
            "Dry-run mode to preview before posting",
            "Reduces manual reconciliation hours each month",
        ],
        how_it_works=[
            "POST /banking/.../ai-reconcile/ sends unmatched lines to Claude",
            "AI returns match suggestions with confidence reasoning",
            "User approves matches — nothing auto-posts without review",
        ],
        example_prompt="Reconcile HDFC account for March — show suggested matches.",
        image=None,
        notes="Month-end close slide for CFO. Human always approves — AI accelerates, doesn't replace judgment.",
        accent=GREEN,
    )

    slide_recruitment_overview(prs, total, next_idx())

    slide_ai_feature(
        prs, total, next_idx(),
        title="JD Generator — professional job posts in seconds",
        tag="AI · Recruitment",
        who="HR admins with recruitment.manage permission",
        bullets=[
            "Input: role title, department, skills, experience, salary band",
            "Output: 350–450 word JD with responsibilities & requirements",
            "Structured JSON for careers page and ATS pipeline",
        ],
        how_it_works=[
            "POST /recruitment/ai/generate-jd/ → Claude with company context",
            "Falls back to template if API key not configured",
            "Tenant-scoped — uses your company name automatically",
        ],
        example_prompt="Generate JD for Senior React Developer, 5 years, Bangalore hybrid.",
        image=GENERATED / "saptta-recruitment-ai-flow.png",
        notes="HR saves 30–60 min per role. Show output quality vs copy-paste from ChatGPT.",
        accent=BLUE,
    )

    slide_ai_feature(
        prs, total, next_idx(),
        title="Resume Parse & Rank — shortlist with evidence",
        tag="AI · Recruitment",
        who="HR recruiters and hiring managers",
        bullets=[
            "Parse PDF resume → skills, experience, education JSON",
            "Rank candidates against job requirements with scores",
            "Explain why candidate A scored higher than B",
        ],
        how_it_works=[
            "Resume parser extracts structured fields via Claude",
            "Ranker compares candidate profile to JD requirements",
            "Pool ranking available for bulk applicant review",
        ],
        example_prompt="Rank these 12 applicants for the DevOps engineer role.",
        image=None,
        notes="Emphasize human makes hire decision. AI removes first-pass screening bias and saves hours.",
        accent=PURPLE,
    )

    slide_ai_feature(
        prs, total, next_idx(),
        title="Offer Letter Generator — consistent, compliant offers",
        tag="AI · Recruitment",
        who="HR admins after candidate selection",
        bullets=[
            "Generate offer letter from approved compensation terms",
            "Uses tenant branding and role details",
            "HR reviews and edits before sending to candidate",
        ],
        how_it_works=[
            "POST /recruitment/ai/generate-offer/ with candidate + CTC details",
            "Claude drafts formal offer with standard clauses",
            "Never auto-sends — always human approval step",
        ],
        example_prompt="Draft offer for Priya Sharma, ₹18 LPA, joining July 1.",
        image=None,
        notes="Close the hire loop. Pair with JD and rank slides for full recruitment story.",
        accent=GREEN,
    )

    slide_ai_feature(
        prs, total, next_idx(),
        title="Performance Review AI — fair drafts from manager notes",
        tag="AI · Performance",
        who="Managers during review cycles",
        bullets=[
            "Manager writes 3–10 bullet notes → structured review draft",
            "Flags biased/coded language (aggressive, emotional, etc.)",
            "Never auto-submits — manager always edits before final",
        ],
        how_it_works=[
            "Heuristic bias term scan before LLM call",
            "Claude structures: strengths, improvements, goals",
            "Prior review context included for continuity",
        ],
        example_prompt="Turn these notes into a review draft for Ankit's Q2 performance.",
        image=None,
        notes="DEI angle for CHRO: catches gendered language managers don't notice. Manager stays accountable.",
        accent=ORANGE,
    )

    slide_ai_feature(
        prs, total, next_idx(),
        title="Saptta Guide — product help without live data",
        tag="AI · Platform",
        who="Prospects, new workspace owners, any user",
        bullets=[
            "Answers what's in each plan, how payroll sync works",
            "No access to tenant financial or HR records",
            "Routes data questions to HR or Finance Assistant",
        ],
        how_it_works=[
            "General AI chat on marketing site and post-login help",
            "Product knowledge base in system prompt",
            "Guest chat for pre-signup visitors",
        ],
        example_prompt="What's included in Complete vs HRMS only? How does payroll sync work?",
        image=GENERATED / "saptta-dual-ai-assistants.png",
        notes="For pre-sales and onboarding. Safe to expose publicly — no PII access.",
        accent=PURPLE,
    )

    slide_modules(prs, total, next_idx())
    slide_compliance(prs, total, next_idx())
    slide_pricing(prs, total, next_idx())
    slide_outcomes(prs, total, next_idx())
    slide_cta(prs, total, next_idx())

    prs.save(OUT)
    print(f"Created: {OUT}")
    print(f"Slides: {total}")
    for p in SHOTS:
        print(f"  [{'OK' if p.is_file() else 'MISSING'}] {p.name}")
    for name in [
        "saptta-ai-three-assistants.png",
        "saptta-payroll-ledger-sync.png",
        "saptta-recruitment-ai-flow.png",
    ]:
        p = GENERATED / name
        print(f"  [{'OK' if p.is_file() else 'MISSING'}] generated/{name}")


if __name__ == "__main__":
    build()
