#!/usr/bin/env python3
"""Generate Saptta executive showcase PowerPoint from slide content."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Saptta-Executive-Showcase-Jun2026.pptx"

BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDES = [
    {
        "title": "One platform for people and money",
        "subtitle": "Saptta HRMS + Saptta Accounts — hire-to-retire HR and GST-ready finance on one login.",
        "bullets": ["2-in-1 HR + Finance", "30 employees included", "$95 Complete / month", "14-day free trial"],
        "title_slide": True,
    },
    {
        "title": "Growing companies lose weeks to disconnected HR & Finance",
        "bullets": [
            "HR runs payroll in one tool; Finance gets another export",
            "Leave balances don't match attendance",
            "Managers email HR — no audit trail",
            "Saptta: one employee record, one truth, manager self-service",
        ],
    },
    {
        "title": "Saptta at a glance",
        "bullets": [
            "Saptta HRMS — $59/mo — people, time, leave, payroll, performance",
            "Saptta Accounts — $59/mo — invoicing, GST, ledger, Tally export",
            "Saptta Complete — $95/mo — both products, save 20%",
        ],
    },
    {
        "title": "Who we serve",
        "bullets": [
            "India SMB: 30–500 employees, full PF/ESI/PT/TDS/Form 16",
            "GCC & Kuwait: core HR live, payroll MVP with PIFSS & indemnity",
            "Industries: IT, manufacturing, retail, professional services",
        ],
    },
    {
        "title": "Hire to retire — one workforce lifecycle",
        "bullets": [
            "Recruit · Onboard · Operate · Pay · Grow · Comply · Exit",
            "ATS, careers pages, AI ranking, hire → employee",
            "Dashboard for company heads: headcount, present/absent, priorities",
        ],
    },
    {
        "title": "Saptta Accounts — books your CA will use",
        "bullets": [
            "Invoicing, purchases, inventory, bank reconciliation",
            "GST returns, ledger, trial balance, Tally XML",
            "Complete plan: payroll journals flow to finance",
        ],
    },
    {
        "title": "Why unified beats fragmented",
        "bullets": [
            "Trust headcount — HR and Finance aligned",
            "Faster month-end — no manual re-keying",
            "Manager self-service — leave, attendance, reviews",
            "50 employees on Complete: $125/month all-in",
        ],
    },
    {
        "title": "India payroll & statutory depth",
        "bullets": [
            "PF ECR · ESI · PT · TDS · Form 16 · Gratuity accrual",
            "Pre-payroll review → run → approve → publish payslips",
            "Honest: TDS simplified — CA validates before 24Q filing",
        ],
    },
    {
        "title": "Kuwait & GCC — core HR + payroll MVP",
        "bullets": [
            "Jurisdiction signup: KW, AE, SA, BH, OM, QA",
            "PIFSS, GOSI, indemnity/EOS accrual, IBAN exports",
            "Validate statutory rates with local payroll consultant",
        ],
    },
    {
        "title": "Outcomes in 90 days",
        "bullets": [
            "60% less payroll prep time",
            "Same-day payslip publish to employees",
            "Zero spreadsheet re-keying between HR and Finance",
            "Full audit trail on sensitive changes",
        ],
    },
    {
        "title": "Simple pricing — under HR-only alternatives",
        "bullets": [
            "HRMS $59 · Accounts $59 · Complete $95 (30 employees)",
            "India: ₹7,999 Complete + GST at checkout",
            "vs Keka ~₹9,999 HR only · vs ZenHR ~$150 HR only",
        ],
    },
    {
        "title": "Next step: 45-minute executive demo",
        "bullets": [
            "Step 1: Discovery — headcount, current tools, pain points",
            "Step 2: Live demo — HR + Finance on one login",
            "Step 3: Trial workspace in 24 hours",
            "hello@saptta.com · saptta.com · 14-day free trial",
        ],
    },
]


def _style_title(shape, *, size=32, color=DARK, bold=True):
    tf = shape.text_frame
    tf.word_wrap = True
    for i, para in enumerate(tf.paragraphs):
        para.font.size = Pt(size if i == 0 else 18)
        para.font.bold = bold if i == 0 else False
        para.font.color.rgb = color
        para.font.name = "Segoe UI"


def _add_bullets(text_frame, items, *, size=16, color=MUTED):
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        is_title = spec.get("title_slide")

        if is_title:
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = DARK
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.2))
            _style_title(title_box, size=40, color=WHITE)
            title_box.text_frame.text = spec["title"]
            sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(10), Inches(1.2))
            sub.text_frame.text = spec.get("subtitle", "")
            for para in sub.text_frame.paragraphs:
                para.font.size = Pt(20)
                para.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
                para.font.name = "Segoe UI"
        else:
            # Header bar
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
            bar.fill.solid()
            bar.fill.fore_color.rgb = BLUE
            bar.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11), Inches(1.2))
            _style_title(title_box, size=28, color=DARK)
            title_box.text_frame.text = spec["title"]

            body = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.5))
            _add_bullets(body.text_frame, spec["bullets"])

        # Footer
        foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(12), Inches(0.4))
        foot.text_frame.text = f"Saptta · Executive Showcase · Slide {idx}/12 · Jun 2026"
        for para in foot.text_frame.paragraphs:
            para.font.size = Pt(10)
            para.font.color.rgb = MUTED if not is_title else RGBColor(0x94, 0xA3, 0xB8)
            para.font.name = "Segoe UI"

    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    build()
