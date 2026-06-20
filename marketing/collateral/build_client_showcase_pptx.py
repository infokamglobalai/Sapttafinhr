#!/usr/bin/env python3
"""
Build one client-facing Saptta showcase PowerPoint with product screenshots,
pricing cards, and polished layout.

Usage:
    python marketing/collateral/build_client_showcase_pptx.py

Output:
    marketing/collateral/Saptta-Client-Showcase.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT = Path(__file__).resolve().parent / "Saptta-Client-Showcase.pptx"
SHOTS = [
    ROOT / "shots" / "02-switcher.png",
    ROOT / "shots" / "04-hr.png",
    ROOT / "shots" / "03-finance.png",
    ROOT / "shots" / "01-superadmin.png",
]

# Brand palette
BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_DARK = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_SOFT = RGBColor(0xEF, 0xF6, 0xFF)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
GREEN_SOFT = RGBColor(0xEC, 0xFD, 0xF5)
ORANGE_SOFT = RGBColor(0xFF, 0xF7, 0xED)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Segoe UI"


def rgb(r, g, b) -> RGBColor:
    return RGBColor(r, g, b)


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(
    shape,
    text: str,
    *,
    size=14,
    color=NAVY,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def _footer(slide, idx: int, total: int):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), H - Inches(0.38), W, Inches(0.38))
    _fill(bar, LIGHT_BG)
    box = slide.shapes.add_textbox(Inches(0.6), H - Inches(0.32), Inches(12), Inches(0.28))
    _text(
        box,
        f"Saptta  ·  Client Showcase  ·  {idx}/{total}  ·  saptta.com  ·  hello@saptta.com",
        size=9,
        color=SLATE,
    )


def _header(slide, title: str, tag: str = ""):
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.35), Inches(0.45), Inches(0.45))
    _fill(logo, BLUE)
    logo.adjustments[0] = 0.15
    lt = slide.shapes.add_textbox(Inches(0.62), Inches(0.42), Inches(0.35), Inches(0.35))
    _text(lt, "S", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    brand = slide.shapes.add_textbox(Inches(1.05), Inches(0.38), Inches(2), Inches(0.4))
    _text(brand, "Saptta", size=14, color=NAVY, bold=True)

    if tag:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(0.78), Inches(1.6), Inches(0.28))
        _fill(pill, BLUE_SOFT)
        pill.line.color.rgb = rgb(0xBF, 0xDB, 0xFE)
        pt = slide.shapes.add_textbox(Inches(1.12), Inches(0.8), Inches(1.5), Inches(0.24))
        _text(pt, tag.upper(), size=8, color=BLUE, bold=True)

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12), Inches(0.75))
    _text(title_box, title, size=28, color=NAVY, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.95), Inches(12.2), Inches(0.03))
    _fill(line, BORDER)


def _card(slide, left, top, width, height, title, lines, *, accent=BLUE, featured=False):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.08
    if featured:
        _fill(rect, BLUE_SOFT)
        rect.line.color.rgb = BLUE
        rect.line.width = Pt(2)
    else:
        _fill(rect, WHITE)
        rect.line.color.rgb = BORDER
        rect.line.width = Pt(1)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    _fill(stripe, accent)

    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.35), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = FONT
    p0.font.size = Pt(12 if not featured else 13)
    p0.font.bold = True
    p0.font.color.rgb = NAVY if not featured else BLUE_DARK
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(10)
        p.font.color.rgb = SLATE
        p.space_before = Pt(4)


def _stat(slide, left, top, value, label, *, primary=False):
    w, h = Inches(2.85), Inches(1.35)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.adjustments[0] = 0.12
    _fill(box, BLUE_SOFT if primary else LIGHT_BG)
    if primary:
        box.line.color.rgb = rgb(0xBF, 0xDB, 0xFE)
    else:
        box.line.color.rgb = BORDER

    v = slide.shapes.add_textbox(left, top + Inches(0.22), w, Inches(0.55))
    _text(v, value, size=26 if primary else 22, color=BLUE if primary else NAVY, bold=True, align=PP_ALIGN.CENTER)
    l = slide.shapes.add_textbox(left, top + Inches(0.78), w, Inches(0.4))
    _text(l, label, size=9, color=SLATE, align=PP_ALIGN.CENTER)


def _image(slide, path: Path, left, top, width, *, caption: str = ""):
    if path.is_file():
        pic = slide.shapes.add_picture(str(path), left, top, width=width)
        # Rounded frame effect — shadow box behind
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left - Inches(0.06),
            top - Inches(0.06),
            width + Inches(0.12),
            pic.height + Inches(0.12),
        )
        frame.adjustments[0] = 0.04
        _fill(frame, WHITE)
        frame.line.color.rgb = BORDER
        # Move picture to front (re-add would be complex; frame under is ok visually via order)
        spTree = slide.shapes._spTree
        sp = pic._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    else:
        placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(4))
        placeholder.adjustments[0] = 0.04
        _fill(placeholder, LIGHT_BG)
        placeholder.line.color.rgb = BORDER
        _text(
            slide.shapes.add_textbox(left, top + Inches(1.5), width, Inches(1)),
            f"[Screenshot: {path.name}]",
            size=11,
            color=SLATE,
            align=PP_ALIGN.CENTER,
        )
    if caption:
        cap = slide.shapes.add_textbox(left, top + Inches(4.05) if path.is_file() else top + Inches(4.15), width, Inches(0.35))
        _text(cap, caption, size=10, color=SLATE, bold=True)


def slide_title(prs, total):
    s = _blank(prs)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    # Accent orb
    orb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.2), Inches(5), Inches(5))
    _fill(orb, BLUE_DARK)
    orb.line.fill.background()

    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.0), Inches(2.4), Inches(0.35))
    _fill(badge, rgb(0x1E, 0x3A, 0x8A))
    b = s.shapes.add_textbox(Inches(1.0), Inches(1.03), Inches(2.2), Inches(0.3))
    _text(b, "CLIENT SHOWCASE · 2026", size=9, color=rgb(0x93, 0xC5, 0xFD), bold=True)

    t = s.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(10), Inches(1.8))
    _text(t, "One platform for\npeople & money", size=44, color=WHITE, bold=True)

    sub = s.shapes.add_textbox(Inches(0.88), Inches(3.35), Inches(9.5), Inches(1))
    _text(
        sub,
        "Saptta HRMS + Saptta Accounts — workforce, payroll, and finance\non a single login. Built for India. Ready for GCC.",
        size=18,
        color=rgb(0xCB, 0xD5, 0xE1),
    )

    stats_y = Inches(4.55)
    for i, (val, lbl) in enumerate([
        ("2-in-1", "HR + Finance"),
        ("30", "Employees incl."),
        ("$95", "Complete / mo"),
        ("14 days", "Free trial"),
    ]):
        _stat(s, Inches(0.85) + i * Inches(3.05), stats_y, val, lbl, primary=(i == 2))

    _footer(s, 1, total)
    return s


def slide_problem(prs, total, idx):
    s = _blank(prs)
    _header(s, "The problem every growing company faces", "Challenge")
    _card(
        s, Inches(0.55), Inches(2.15), Inches(5.9), Inches(3.8),
        "Without Saptta",
        [
            "✗  HR and Finance use different tools",
            "✗  Payroll re-keyed into accounting",
            "✗  Leave ≠ attendance in reports",
            "✗  Managers email HR for approvals",
            "✗  No audit trail on sensitive data",
        ],
        accent=ORANGE,
    )
    _card(
        s, Inches(6.85), Inches(2.15), Inches(5.9), Inches(3.8),
        "With Saptta Complete",
        [
            "✓  One employee record — one truth",
            "✓  Attendance → payroll → payslip flow",
            "✓  Manager & employee self-service",
            "✓  PF, GST, Tally exports built-in",
            "✓  Encrypted PII + role-based access",
        ],
        accent=GREEN,
        featured=True,
    )
    callout = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.55))
    callout.adjustments[0] = 0.2
    _fill(callout, BLUE_SOFT)
    callout.line.color.rgb = rgb(0xBF, 0xDB, 0xFE)
    c = s.shapes.add_textbox(Inches(0.75), Inches(6.12), Inches(11.8), Inches(0.45))
    _text(c, "CEO insight: the hidden cost is leadership time reconciling two versions of your own company.", size=11, color=BLUE_DARK, bold=True)
    _footer(s, idx, total)


def slide_platform(prs, total, idx):
    s = _blank(prs)
    _header(s, "Two products. One platform.", "Platform")
    plans = [
        ("Saptta HRMS", "$59/mo", ["Payroll & statutory", "Leave & attendance", "Recruitment & performance", "Manager + employee ESS"], BLUE, False),
        ("Saptta Accounts", "$59/mo", ["Invoicing & GST", "Ledger & bank rec", "Inventory & purchases", "Unlimited finance users"], rgb(0x7C, 0x3A, 0xED), False),
        ("Saptta Complete ★", "$95/mo", ["Everything in HRMS + Accounts", "20% savings vs separate", "30 employees included", "Recommended for owners"], BLUE, True),
    ]
    x0 = Inches(0.55)
    for i, (name, price, feats, accent, feat) in enumerate(plans):
        left = x0 + i * Inches(4.15)
        _card(s, left, Inches(2.2), Inches(3.95), Inches(3.5), name, [price, ""] + feats, accent=accent, featured=feat)
    flow = s.shapes.add_textbox(Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.45))
    _text(
        flow,
        "Sign up  →  Product switcher  →  HR workspace  ↔  Finance workspace  →  Single SSO",
        size=12,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )
    _footer(s, idx, total)


def slide_screenshot(prs, total, idx, title, tag, img_path, caption, bullets):
    s = _blank(prs)
    _header(s, title, tag)
    _image(s, img_path, Inches(6.6), Inches(2.15), Inches(6.2), caption=caption)
    tb = s.shapes.add_textbox(Inches(0.55), Inches(2.2), Inches(5.8), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.name = FONT
        p.font.size = Pt(13)
        p.font.color.rgb = SLATE if i > 0 else NAVY
        p.font.bold = i == 0 and False
        p.space_after = Pt(10)
    _footer(s, idx, total)


def slide_modules(prs, total, idx):
    s = _blank(prs)
    _header(s, "Hire to retire — full HR lifecycle", "HRMS")
    modules = [
        ("Recruit", "ATS, careers, AI rank"),
        ("Onboard", "Checklists, documents"),
        ("Operate", "Attendance, leave"),
        ("Pay", "Run, payslip, bank"),
        ("Grow", "Performance reviews"),
        ("Comply", "Letters, policies"),
        ("Engage", "ESS portal"),
        ("Exit", "Settlement, revoke"),
    ]
    for i, (name, desc) in enumerate(modules):
        col, row = i % 4, i // 4
        left = Inches(0.55) + col * Inches(3.1)
        top = Inches(2.15) + row * Inches(1.85)
        _card(s, left, top, Inches(2.85), Inches(1.55), name, [desc], accent=BLUE)
    _footer(s, idx, total)


def slide_compliance(prs, total, idx):
    s = _blank(prs)
    _header(s, "Compliance built for real businesses", "Trust")
    _card(
        s, Inches(0.55), Inches(2.15), Inches(5.9), Inches(4.2),
        "India (production-ready)",
        ["PF ECR · ESI · PT · TDS · Form 16", "Gratuity accrual + exit estimate", "Salary register · Bank advice · Tally XML", "Pre-payroll review → approve → publish"],
        accent=GREEN,
        featured=True,
    )
    _card(
        s, Inches(6.85), Inches(2.15), Inches(5.9), Inches(4.2),
        "GCC & Kuwait (demo-ready MVP)",
        ["KW · AE · SA jurisdiction signup", "PIFSS · GOSI · indemnity accrual", "IBAN · WPS-style exports", "Validate rates with local consultant"],
        accent=ORANGE,
    )
    _footer(s, idx, total)


def slide_pricing(prs, total, idx):
    s = _blank(prs)
    _header(s, "Simple pricing — beat HR-only tools", "Pricing")
    rows = [
        ["Plan", "USD / mo", "India INR", "Best for"],
        ["HRMS", "$59", "₹4,999", "HR team replacing legacy HR"],
        ["Accounts", "$59", "₹4,999", "Finance + CA workflow"],
        ["Complete ★", "$95", "₹7,999", "Owners wanting one platform"],
    ]
    tbl = s.shapes.add_table(4, 4, Inches(0.55), Inches(2.15), Inches(7.5), Inches(2.2)).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.name = FONT
                para.font.size = Pt(10 if r else 9)
                para.font.bold = r == 0 or (r == 3 and c == 0)
                para.font.color.rgb = NAVY if r == 3 else (SLATE if r else WHITE)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif r == 3:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_SOFT

    _card(
        s, Inches(8.35), Inches(2.15), Inches(4.4), Inches(2.2),
        "At 30 employees",
        ["Saptta Complete: $95", "Keka HR only: ~₹9,999", "ZenHR GCC: ~$150", "You save 20%+ vs separate"],
        accent=BLUE,
        featured=True,
    )
    for i, (val, lbl) in enumerate([("50", "emp → $125"), ("100", "emp → $200"), ("+30", "+$1.50/emp")]):
        _stat(s, Inches(8.35) + i * Inches(1.5), Inches(4.55), val, lbl)
    _footer(s, idx, total)


def slide_outcomes(prs, total, idx):
    s = _blank(prs)
    _header(s, "Results company heads see in 90 days", "Outcomes")
    for i, (val, lbl) in enumerate([
        ("60%", "Less payroll prep"),
        ("Same day", "Payslip publish"),
        ("Zero", "Spreadsheet re-key"),
        ("100%", "Audit on changes"),
    ]):
        _stat(s, Inches(0.55) + i * Inches(3.15), Inches(2.3), val, lbl, primary=(i == 0))
    _card(s, Inches(0.55), Inches(4.0), Inches(5.9), Inches(2.3), "CHRO wins", ["Manager leave approvals", "Recruitment pipeline visible", "Compliance expiry alerts"], accent=BLUE)
    _card(s, Inches(6.85), Inches(4.0), Inches(5.9), Inches(2.3), "CFO wins", ["Payroll cost in finance view", "GST + ledger same platform", "Export-ready for CA"], accent=GREEN)
    _footer(s, idx, total)


def slide_cta(prs, total, idx):
    s = _blank(prs)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    t = s.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(11), Inches(1.2))
    _text(t, "See your company on Saptta", size=36, color=WHITE, bold=True)

    sub = s.shapes.add_textbox(Inches(0.92), Inches(2.35), Inches(10), Inches(0.6))
    _text(sub, "45-minute live demo on your scenarios · Trial workspace in 24 hours", size=16, color=rgb(0x93, 0xC5, 0xFD))

    steps = [
        ("1", "Discovery", "Headcount, tools, pain points"),
        ("2", "Live demo", "HR + Finance, one login"),
        ("3", "Trial", "Your team tests in 24h"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        left = Inches(0.9) + i * Inches(4.05)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.2), Inches(3.7), Inches(1.65))
        box.adjustments[0] = 0.1
        _fill(box, rgb(0x1E, 0x29, 0x3B))
        box.line.color.rgb = rgb(0x33, 0x41, 0x55)
        n = s.shapes.add_textbox(left + Inches(0.2), Inches(3.35), Inches(0.5), Inches(0.4))
        _text(n, num, size=20, color=BLUE, bold=True)
        tt = s.shapes.add_textbox(left + Inches(0.55), Inches(3.32), Inches(2.9), Inches(0.4))
        _text(tt, title, size=14, color=WHITE, bold=True)
        dd = s.shapes.add_textbox(left + Inches(0.2), Inches(3.85), Inches(3.3), Inches(0.7))
        _text(dd, desc, size=11, color=rgb(0x94, 0xA3, 0xB8))

    contact = s.shapes.add_textbox(Inches(0.9), Inches(5.35), Inches(11), Inches(0.8))
    _text(contact, "hello@saptta.com  ·  saptta.com  ·  14-day free trial  ·  Free Complete setup (launch offer)", size=14, color=WHITE, align=PP_ALIGN.CENTER)

    _footer(s, idx, total)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = 11

    slide_title(prs, total)
    slide_problem(prs, total, 2)
    slide_platform(prs, total, 3)
    slide_screenshot(
        prs, total, 4,
        "One login. Two products.",
        "Live product",
        SHOTS[0],
        "Product switcher — HR & Finance",
        [
            "Unified platform login",
            "Switch HR ↔ Finance instantly",
            "Same workspace, same users",
            "No duplicate employee master",
        ],
    )
    slide_screenshot(
        prs, total, 5,
        "HR dashboard built for leaders",
        "Saptta HRMS",
        SHOTS[1],
        "Headcount, attendance, priorities",
        [
            "Today's present / absent panel",
            "Pending leave & payroll actions",
            "22+ employees demo-ready",
            "Manager & employee ESS",
        ],
    )
    slide_screenshot(
        prs, total, 6,
        "Finance your CA will use",
        "Saptta Accounts",
        SHOTS[2],
        "Invoicing, ledger, GST",
        [
            "Full accounting workspace",
            "GST-ready compliance",
            "Tally XML export",
            "Unlimited finance users (flat fee)",
        ],
    )
    slide_modules(prs, total, 7)
    slide_compliance(prs, total, 8)
    slide_pricing(prs, total, 9)
    slide_outcomes(prs, total, 10)
    slide_cta(prs, total, 11)

    prs.save(OUT)
    print(f"Created: {OUT}")
    for p in SHOTS:
        status = "OK" if p.is_file() else "MISSING"
        print(f"  [{status}] {p.name}")


if __name__ == "__main__":
    build()
